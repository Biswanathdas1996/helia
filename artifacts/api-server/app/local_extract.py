from __future__ import annotations

from io import BytesIO


def extract_text_locally(mime_type: str, raw: bytes) -> str | None:
    if mime_type == "application/pdf":
        return _extract_pdf(raw)
    if mime_type in {
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword",
    }:
        return _extract_doc(raw)
    if mime_type in {
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint",
    }:
        return _extract_ppt(raw)
    if mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet":
        return _extract_xlsx(raw)
    return None


def _extract_pdf(raw: bytes) -> str:
    try:
        from pypdf import PdfReader  # type: ignore[import-not-found]
    except Exception as err:  # pragma: no cover
        raise RuntimeError("PDF parser is not installed") from err

    reader = PdfReader(BytesIO(raw))
    parts: list[str] = []
    for page in reader.pages:
        parts.append((page.extract_text() or "").strip())
    return "\n\n".join(p for p in parts if p)


def _extract_doc(raw: bytes) -> str:
    if raw.startswith(b"PK"):
        try:
            from docx import Document  # type: ignore[import-not-found]
        except Exception as err:  # pragma: no cover
            raise RuntimeError("DOCX parser is not installed") from err

        doc = Document(BytesIO(raw))
        parts: list[str] = []
        for p in doc.paragraphs:
            txt = (p.text or "").strip()
            if txt:
                parts.append(txt)
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells if c.text and c.text.strip()]
                if cells:
                    parts.append(" | ".join(cells))
        return "\n".join(parts)

    # Legacy .doc is not supported by python-docx. Let caller fallback to AI path.
    raise RuntimeError("Legacy .doc format requires AI extraction")


def _extract_ppt(raw: bytes) -> str:
    if not raw.startswith(b"PK"):
        raise RuntimeError("Legacy .ppt format requires AI extraction")

    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except Exception as err:  # pragma: no cover
        raise RuntimeError("PPTX parser is not installed") from err

    prs = Presentation(BytesIO(raw))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts.append(f"[Slide {i}]")
        for shape in slide.shapes:
            txt = getattr(shape, "text", "")
            txt = (txt or "").strip()
            if txt:
                parts.append(txt)
    return "\n".join(parts)


def _extract_xlsx(raw: bytes) -> str:
    try:
        from openpyxl import load_workbook  # type: ignore[import-not-found]
    except Exception as err:  # pragma: no cover
        raise RuntimeError("XLSX parser is not installed") from err

    wb = load_workbook(BytesIO(raw), read_only=True, data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"[Sheet: {ws.title}]")
        for row in ws.iter_rows(values_only=True):
            cells = [str(v).strip() for v in row if v is not None and str(v).strip()]
            if cells:
                parts.append(" | ".join(cells))
    return "\n".join(parts)
